"""
AI 最前线 PPT 生成脚本 v2
聚焦最新 AI 热点新闻 + 深度洞察，适合课外分享。
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ──────────────── 设计常量 ────────────────
BG_DARK = RGBColor(0x0F, 0x0F, 0x1A)
BG_CARD = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD_ALT = RGBColor(0x1E, 0x1E, 0x36)
ACCENT_BLUE = RGBColor(0x38, 0xBD, 0xF8)
ACCENT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
ACCENT_GREEN = RGBColor(0x34, 0xD3, 0x99)
ACCENT_ORANGE = RGBColor(0xFB, 0xBF, 0x24)
ACCENT_RED = RGBColor(0xF8, 0x71, 0x71)
ACCENT_PINK = RGBColor(0xF0, 0x6E, 0xCD)
ACCENT_CYAN = RGBColor(0x22, 0xD3, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xC0, 0xC0, 0xD0)
MEDIUM_GRAY = RGBColor(0x78, 0x78, 0x90)
DIM_GRAY = RGBColor(0x50, 0x50, 0x68)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
FONT_BODY = "PingFang SC"

TOTAL_SLIDES = 15


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_line(slide, left, top, width, color):
    return add_rect(slide, left, top, width, Pt(2.5), color)


def add_text(slide, left, top, width, height, text, size=18,
             color=WHITE, bold=False, align=PP_ALIGN.LEFT, spacing=1.5):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT_BODY
    p.alignment = align
    p.line_spacing = Pt(size * spacing)
    return box


def add_lines(slide, left, top, width, height, lines, size=16,
              color=WHITE, bold=False, align=PP_ALIGN.LEFT, spacing=1.5):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = FONT_BODY
        p.alignment = align
        p.space_after = Pt(size * 0.4)
        p.line_spacing = Pt(size * spacing)
    return box


def page_num(slide, n):
    add_text(slide, SLIDE_WIDTH - Inches(1), SLIDE_HEIGHT - Inches(0.45),
             Inches(0.8), Inches(0.35), f"{n}/{TOTAL_SLIDES}",
             size=9, color=DIM_GRAY, align=PP_ALIGN.RIGHT)


def tag(slide, left, top, text, color):
    add_rounded_rect(slide, left, top, Pt(len(text) * 11 + 24), Pt(26), color)
    add_text(slide, left + Pt(12), top + Pt(2), Pt(len(text) * 11), Pt(22),
             text, size=10, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)


# ──────────────── 构建 PPT ────────────────
prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT
blank = prs.slide_layouts[6]

# ============================================================
# 1. 封面
# ============================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)

# 装饰光晕
add_circle(s, Inches(-2), Inches(-2), Inches(6), RGBColor(0x0A, 0x1A, 0x3A))
add_circle(s, Inches(9), Inches(4), Inches(5), RGBColor(0x1A, 0x0A, 0x2A))

add_text(s, Inches(1.2), Inches(1.2), Inches(11), Inches(0.4),
         "AI 最前线  ·  2026 年 5 月", size=16, color=ACCENT_BLUE, bold=True)

add_text(s, Inches(1.2), Inches(2.0), Inches(11), Inches(1.5),
         "这周，AI 世界\n发生了什么疯狂的事？",
         size=52, color=WHITE, bold=True, spacing=1.3)

add_line(s, Inches(1.2), Inches(4.3), Inches(2), ACCENT_BLUE)

add_lines(s, Inches(1.2), Inches(4.8), Inches(8), Inches(1.5),
          ["OpenAI 推翻 80 年数学猜想 · Google 发布 Gemini 3.5 · 国产模型全面崛起",
           "AI Agent 开始替人打工 · 强化学习让 AI 自己「顿悟」"],
          size=16, color=MEDIUM_GRAY, spacing=1.8)

add_text(s, Inches(1.2), Inches(6.2), Inches(8), Inches(0.4),
         "一份来自 AI 学习小组的课外见闻分享", size=14, color=DIM_GRAY)

# ============================================================
# 2. 本期看点
# ============================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
page_num(s, 2)

add_text(s, Inches(1), Inches(0.6), Inches(3), Inches(0.4),
         "📌 本期看点", size=14, color=ACCENT_CYAN, bold=True)
add_text(s, Inches(1), Inches(1.1), Inches(11), Inches(0.7),
         "5 个值得你知道的 AI 大事件", size=38, color=WHITE, bold=True)

items = [
    ("🔥", "01", "AI 推翻 80 年数学猜想", "2天前｜OpenAI 的模型独立证明了一个数学家做不到的事", ACCENT_RED),
    ("🚀", "02", "Google I/O 大战全面开打", "本周｜Gemini 3.5 · Spark Agent · Omni 世界模型", ACCENT_BLUE),
    ("🇨🇳", "03", "国产大模型「卷疯了」", "本周｜阿里千问 · 蚂蚁百灵万亿参数开源", ACCENT_GREEN),
    ("🧠", "04", "DeepSeek 的「顿悟」实验", "回顾｜AI 如何用强化学习自己学会思考", ACCENT_PURPLE),
    ("🤖", "05", "AI Agent 时代来了", "趋势｜Manus · Gemini Spark · 从聊天到干活", ACCENT_ORANGE),
]

for i, (emoji, num, title, desc, color) in enumerate(items):
    y = Inches(2.2) + Inches(i * 0.98)
    add_rounded_rect(s, Inches(1), y, Inches(11.3), Inches(0.85), BG_CARD)
    add_text(s, Inches(1.3), y + Inches(0.12), Inches(0.5), Inches(0.5),
             emoji, size=22)
    add_text(s, Inches(1.9), y + Inches(0.15), Inches(0.5), Inches(0.5),
             num, size=18, color=color, bold=True)
    add_text(s, Inches(2.5), y + Inches(0.15), Inches(4), Inches(0.5),
             title, size=18, color=WHITE, bold=True)
    add_text(s, Inches(6.8), y + Inches(0.18), Inches(5.3), Inches(0.5),
             desc, size=13, color=MEDIUM_GRAY)

# ============================================================
# 3. OpenAI 推翻数学猜想
# ============================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
page_num(s, 3)

tag(s, Inches(1), Inches(0.6), "HOT", ACCENT_RED)
add_text(s, Inches(1.8), Inches(0.55), Inches(3), Inches(0.4),
         "5月20日  ·  2天前", size=12, color=MEDIUM_GRAY)

add_text(s, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         "AI 独立推翻了一个 80 年的数学猜想 🤯", size=36, color=WHITE, bold=True)
add_line(s, Inches(1), Inches(2.1), Inches(1.5), ACCENT_RED)

# 左：事件
add_rounded_rect(s, Inches(1), Inches(2.5), Inches(5.8), Inches(4.5), BG_CARD)
add_text(s, Inches(1.3), Inches(2.7), Inches(5.2), Inches(0.4),
         "📖 发生了什么？", size=17, color=ACCENT_RED, bold=True)
add_lines(s, Inches(1.3), Inches(3.3), Inches(5.2), Inches(3.5),
          ["1946 年，数学家 Erdős 提出一个几何猜想：",
           "「在平面上放 n 个点，最多能有多少对点距离恰好为 1？」",
           "",
           "80 年来，数学家们一直认为正方形网格是最优的。",
           "",
           "然后上周二，OpenAI 的推理模型自己发现了",
           "一族全新的构造方法，直接推翻了这个猜想。",
           "",
           "菲尔兹奖得主 Tim Gowers 亲自验证：",
           "\"这个证明是正确的，而且非常巧妙。\""],
          size=14, color=LIGHT_GRAY, spacing=1.4)

# 右：为什么重要
add_rounded_rect(s, Inches(7.2), Inches(2.5), Inches(5.3), Inches(4.5), BG_CARD)
add_text(s, Inches(7.5), Inches(2.7), Inches(4.7), Inches(0.4),
         "💡 为什么这件事炸裂？", size=17, color=ACCENT_ORANGE, bold=True)
add_lines(s, Inches(7.5), Inches(3.3), Inches(4.7), Inches(3.5),
          ["这是 AI 第一次独立解决一个",
           "数学界公认的重要开放问题。",
           "",
           "注意关键词：「独立」。",
           "不是人类指导它怎么做，",
           "是 AI 自己找到了一条",
           "80 年来没人想到的路。",
           "",
           "从 AlphaGo 的「第37手」到这次，",
           "AI 正在从「模仿人类」走向「超越人类」。",
           "",
           "🔮 这意味着 AI 不只是工具，",
           "它正在变成一个「合作者」。"],
          size=14, color=LIGHT_GRAY, spacing=1.4)

# ============================================================
# 4. Google I/O 2026
# ============================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
page_num(s, 4)

tag(s, Inches(1), Inches(0.6), "本周", ACCENT_BLUE)
add_text(s, Inches(1.8), Inches(0.55), Inches(5), Inches(0.4),
         "Google I/O 2026  ·  5月20日", size=12, color=MEDIUM_GRAY)
add_text(s, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         "Google 的反击：三大王炸一起甩 💣", size=36, color=WHITE, bold=True)
add_line(s, Inches(1), Inches(2.1), Inches(1.5), ACCENT_BLUE)

cards = [
    ("Gemini 3.5 Flash",
     "🏎️ 更快更便宜",
     ["速度是同类产品的 4 倍",
      "价格只有 1/3",
      "超越上一代旗舰 Gemini 3.1 Pro",
      "月活 9 亿用户（同比翻倍）",
      "",
      "皮查伊：大客户每年能省 10 亿美元"],
     ACCENT_BLUE),
    ("Gemini Spark",
     "🤖 24h 私人 AI 管家",
     ["不只是聊天，能直接帮你干活",
      "深度打通 Gmail、日历、文档",
      "能跨 App 推理和执行任务",
      "24 小时不间断运行",
      "",
      "你的第一个「AI 员工」？"],
     ACCENT_PURPLE),
    ("Omni 世界模型",
     "🌍 模拟真实世界",
     ["不是聊天模型，是「物理引擎」",
      "能预测物理世界的结果",
      "支持图片/音频/视频编辑",
      "原生多模态输入输出",
      "",
      "AI 开始「理解」物理世界了"],
     ACCENT_GREEN),
]

for i, (name, subtitle, lines, color) in enumerate(cards):
    x = Inches(1) + Inches(i * 3.9)
    add_rounded_rect(s, x, Inches(2.5), Inches(3.6), Inches(4.5), BG_CARD)
    add_line(s, x + Inches(0.3), Inches(2.65), Inches(3), color)
    add_text(s, x + Inches(0.3), Inches(2.85), Inches(3), Inches(0.5),
             name, size=22, color=color, bold=True)
    add_text(s, x + Inches(0.3), Inches(3.35), Inches(3), Inches(0.4),
             subtitle, size=14, color=LIGHT_GRAY)
    add_lines(s, x + Inches(0.3), Inches(3.9), Inches(3), Inches(2.8),
              lines, size=13, color=MEDIUM_GRAY, spacing=1.4)

# ============================================================
# 5. 国产大模型
# ============================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
page_num(s, 5)

tag(s, Inches(1), Inches(0.6), "国产", ACCENT_GREEN)
add_text(s, Inches(1.8), Inches(0.55), Inches(5), Inches(0.4),
         "本周动态", size=12, color=MEDIUM_GRAY)
add_text(s, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         "国产大模型这周「卷」到了什么程度？🔥", size=36, color=WHITE, bold=True)
add_line(s, Inches(1), Inches(2.1), Inches(1.5), ACCENT_GREEN)

# 阿里
add_rounded_rect(s, Inches(1), Inches(2.5), Inches(5.8), Inches(2.2), BG_CARD)
add_text(s, Inches(1.3), Inches(2.7), Inches(5.2), Inches(0.4),
         "🏆 阿里千问 Qwen3.7-Max", size=18, color=ACCENT_BLUE, bold=True)
add_lines(s, Inches(1.3), Inches(3.2), Inches(5.2), Inches(1.3),
          ["国产模型排名第一 · 编程/推理多项指标领先",
           "能独立完成 35 小时的硬件优化任务",
           "一个 AI 干了工程师好几天的活"],
          size=14, color=LIGHT_GRAY, spacing=1.5)

# 蚂蚁
add_rounded_rect(s, Inches(7.2), Inches(2.5), Inches(5.3), Inches(2.2), BG_CARD)
add_text(s, Inches(7.5), Inches(2.7), Inches(4.7), Inches(0.4),
         "🐜 蚂蚁百灵 Ling-2.6-1T", size=18, color=ACCENT_GREEN, bold=True)
add_lines(s, Inches(7.5), Inches(3.2), Inches(4.7), Inches(1.3),
          ["万亿参数  ·  完全开源",
           "综合智能对标 GPT-5.4",
           "「快思考」机制，推理效率拉满"],
          size=14, color=LIGHT_GRAY, spacing=1.5)

# 大背景
add_rounded_rect(s, Inches(1), Inches(5.0), Inches(11.5), Inches(2.1), BG_CARD)
add_text(s, Inches(1.3), Inches(5.2), Inches(10.9), Inches(0.4),
         "📊 大背景：中国 AI 的「iPhone 时刻」", size=18, color=ACCENT_ORANGE, bold=True)
add_lines(s, Inches(1.3), Inches(5.7), Inches(10.9), Inches(1.2),
          ["半年前 DeepSeek R1 震惊全球（性能对标 OpenAI o1，成本仅 3-5%），打破了「只有硅谷能做顶级AI」的神话",
           "现在阿里、蚂蚁、字节等公司全面跟进，开源模型质量已经不输闭源大厂",
           "结论：AI 不是美国的专利，中国在这场竞赛里已经是第一梯队的选手"],
          size=14, color=LIGHT_GRAY, spacing=1.6)

# ============================================================
# 6. DeepSeek R1 故事
# ============================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
page_num(s, 6)

tag(s, Inches(1), Inches(0.6), "深度", ACCENT_PURPLE)
add_text(s, Inches(1.8), Inches(0.55), Inches(5), Inches(0.4),
         "最值得了解的技术突破", size=12, color=MEDIUM_GRAY)
add_text(s, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         "DeepSeek R1：AI 是怎么自己「顿悟」的？🧠", size=36, color=WHITE, bold=True)
add_line(s, Inches(1), Inches(2.1), Inches(1.5), ACCENT_PURPLE)

# 左侧故事
add_rounded_rect(s, Inches(1), Inches(2.5), Inches(6), Inches(4.5), BG_CARD)
add_text(s, Inches(1.3), Inches(2.7), Inches(5.4), Inches(0.4),
         "📖 一个让全世界震惊的实验", size=17, color=ACCENT_PURPLE, bold=True)
add_lines(s, Inches(1.3), Inches(3.3), Inches(5.4), Inches(3.5),
          ["传统做法：先教 AI 模仿人类的解题步骤",
           "（就像抄学霸的作业）",
           "",
           "DeepSeek 的做法：",
           "只告诉 AI 正确答案是什么，",
           "让它自己想办法怎么解出来。",
           "",
           "结果 AI 自己发展出了：",
           "• 自我反思（\"等等，我算错了\"）",
           "• 多角度思考（\"换个方法试试\"）",
           "• 回溯纠错（\"退回去重来\"）",
           "",
           "这些能力不是人教的，是它自己「悟」出来的。"],
          size=14, color=LIGHT_GRAY, spacing=1.35)

# 右侧数据
add_rounded_rect(s, Inches(7.3), Inches(2.5), Inches(5.2), Inches(2), BG_CARD)
add_text(s, Inches(7.6), Inches(2.7), Inches(4.6), Inches(0.4),
         "📊 硬核数据", size=17, color=ACCENT_CYAN, bold=True)
add_lines(s, Inches(7.6), Inches(3.2), Inches(4.6), Inches(1.2),
          ["AIME 数学测试：15.6% → 71.0%（纯 RL）",
           "性能对标 OpenAI o1",
           "成本仅为 OpenAI 的 3-5% 💰"],
          size=14, color=LIGHT_GRAY, spacing=1.5)

add_rounded_rect(s, Inches(7.3), Inches(4.8), Inches(5.2), Inches(2.2), BG_CARD)
add_text(s, Inches(7.6), Inches(5.0), Inches(4.6), Inches(0.4),
         "🤔 为什么这很重要？", size=17, color=ACCENT_ORANGE, bold=True)
add_lines(s, Inches(7.6), Inches(5.5), Inches(4.6), Inches(1.3),
          ["这证明 AI 可以不靠人类指导，",
           "自己学会复杂推理。",
           "",
           "就像 AlphaGo 的「第37手」——",
           "AI 走出了人类没想到的路。"],
          size=14, color=LIGHT_GRAY, spacing=1.4)

# ============================================================
# 7. 强化学习的 "顿悟" 现象
# ============================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
page_num(s, 7)

add_text(s, Inches(1), Inches(0.6), Inches(3), Inches(0.4),
         "🔬 深入一层", size=14, color=ACCENT_PURPLE, bold=True)
add_text(s, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         "从 AlphaGo 到 ChatGPT：「第37手」现象", size=36, color=WHITE, bold=True)
add_line(s, Inches(1), Inches(2.1), Inches(1.5), ACCENT_PURPLE)

# AlphaGo
add_rounded_rect(s, Inches(1), Inches(2.5), Inches(5.8), Inches(4.3), BG_CARD)
add_text(s, Inches(1.3), Inches(2.7), Inches(5.2), Inches(0.4),
         "♟️ 2016：AlphaGo vs 李世石", size=17, color=ACCENT_BLUE, bold=True)
add_lines(s, Inches(1.3), Inches(3.3), Inches(5.2), Inches(3.3),
          ["第二局，第37手。",
           "",
           "AlphaGo 下了一步所有人类专家",
           "都认为「不可能是好棋」的棋。",
           "",
           "人类专家评估这步棋被选中的概率：",
           "1/10,000",
           "",
           "结果？这步棋直接奠定了胜局。",
           "",
           "赛后分析：这步棋利用了一种",
           "人类几千年围棋史从未发现的策略。"],
          size=14, color=LIGHT_GRAY, spacing=1.35)

# ChatGPT
add_rounded_rect(s, Inches(7.2), Inches(2.5), Inches(5.3), Inches(4.3), BG_CARD)
add_text(s, Inches(7.5), Inches(2.7), Inches(4.7), Inches(0.4),
         "💬 2025-2026：ChatGPT 的同款时刻", size=17, color=ACCENT_GREEN, bold=True)
add_lines(s, Inches(7.5), Inches(3.3), Inches(4.7), Inches(3.3),
          ["同样的事情正在语言模型上重演：",
           "",
           "DeepSeek R1 / OpenAI o3 这些「思考模型」",
           "通过强化学习，自己发展出了",
           "人类从未设计过的推理策略。",
           "",
           "最新证据：OpenAI 的模型刚刚",
           "独立推翻了一个 80 年的数学猜想。",
           "",
           "🔮 未来可能出现：",
           "• AI 发明自己的「思维语言」",
           "• AI 发现人类未知的科学规律",
           "• AI 成为真正的「研究合作者」"],
          size=14, color=LIGHT_GRAY, spacing=1.35)

# ============================================================
# 8. AI Agent 时代
# ============================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
page_num(s, 8)

tag(s, Inches(1), Inches(0.6), "趋势", ACCENT_ORANGE)
add_text(s, Inches(1.8), Inches(0.55), Inches(5), Inches(0.4),
         "2026 年最大的风口", size=12, color=MEDIUM_GRAY)
add_text(s, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         "AI Agent：从「聊天」到「干活」的质变 🤖", size=36, color=WHITE, bold=True)
add_line(s, Inches(1), Inches(2.1), Inches(1.5), ACCENT_ORANGE)

# 对比
add_rounded_rect(s, Inches(1), Inches(2.5), Inches(5.2), Inches(2.2), BG_CARD)
add_text(s, Inches(1.3), Inches(2.7), Inches(4.6), Inches(0.4),
         "😴 以前的 AI（聊天模式）", size=16, color=MEDIUM_GRAY, bold=True)
add_lines(s, Inches(1.3), Inches(3.2), Inches(4.6), Inches(1.3),
          ["你：帮我订明天去上海的高铁票",
           "AI：好的，你可以打开 12306 App...",
           "你：......我知道，我要你帮我订啊！",
           "AI：抱歉，我无法执行实际操作 🤷"],
          size=13, color=LIGHT_GRAY, spacing=1.5)

add_rounded_rect(s, Inches(7.3), Inches(2.5), Inches(5.2), Inches(2.2), BG_CARD)
add_text(s, Inches(7.6), Inches(2.7), Inches(4.6), Inches(0.4),
         "🔥 现在的 AI（Agent 模式）", size=16, color=ACCENT_ORANGE, bold=True)
add_lines(s, Inches(7.6), Inches(3.2), Inches(4.6), Inches(1.3),
          ["你：帮我订明天去上海的高铁票",
           "AI：收到，正在打开 12306...",
           "AI：已选择 G7 次 08:00，二等座 ¥553",
           "AI：需要确认支付吗？✅"],
          size=13, color=LIGHT_GRAY, spacing=1.5)

# Agent 玩家
add_rounded_rect(s, Inches(1), Inches(5.0), Inches(11.5), Inches(2.1), BG_CARD)
add_text(s, Inches(1.3), Inches(5.2), Inches(10.9), Inches(0.4),
         "🏟️ 当前的主要 Agent 玩家", size=17, color=ACCENT_CYAN, bold=True)
add_lines(s, Inches(1.3), Inches(5.7), Inches(10.9), Inches(1.2),
          ["Manus — GAIA 基准 86.5%，远超 OpenAI Deep Research（47.6%），估值 $20 亿，4 个月营收 $9000 万",
           "Gemini Spark（Google）— 打通 Gmail/日历/文档，24h 自动运行，这周刚发布",
           "OpenAI Operator — 能控制浏览器键盘鼠标 · 阿里千问 — 独立完成 35 小时硬件任务"],
          size=13, color=LIGHT_GRAY, spacing=1.6)

# ============================================================
# 9. AI 的阿喀琉斯之踵（有趣洞察）
# ============================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
page_num(s, 9)

add_text(s, Inches(1), Inches(0.6), Inches(3), Inches(0.4),
         "🧩 冷知识", size=14, color=ACCENT_PINK, bold=True)
add_text(s, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         "AI 的「阿喀琉斯之踵」：那些诡异的翻车 🎪", size=36, color=WHITE, bold=True)
add_line(s, Inches(1), Inches(2.1), Inches(1.5), ACCENT_PINK)

# 9.11 vs 9.9
add_rounded_rect(s, Inches(1), Inches(2.5), Inches(5.8), Inches(4.3), BG_CARD)
add_text(s, Inches(1.3), Inches(2.7), Inches(5.2), Inches(0.4),
         "📖 「圣经假说」：9.11 vs 9.9", size=17, color=ACCENT_RED, bold=True)
add_lines(s, Inches(1.3), Inches(3.3), Inches(5.2), Inches(3.3),
          ["问 AI：9.11 和 9.9 哪个大？",
           "AI 有时候说 9.11 更大 ❌",
           "",
           "这怎么可能？AI 能做微积分却不会比大小？",
           "",
           "有一个未被证实但很有趣的假说：",
           "9.11 和 9.9 长得像圣经章节编号。",
           "（圣经第9章第11节 确实在 第9节 后面）",
           "",
           "AI 大脑里负责处理圣经内容的神经元",
           "被意外激活了，干扰了数学比较。",
           "",
           "这说明 AI 不是「理解」数字，",
           "而是在「模式匹配」，模式错了就翻车。"],
          size=13, color=LIGHT_GRAY, spacing=1.35)

# 瑞士奶酪
add_rounded_rect(s, Inches(7.2), Inches(2.5), Inches(5.3), Inches(4.3), BG_CARD)
add_text(s, Inches(7.5), Inches(2.7), Inches(4.7), Inches(0.4),
         "🧀 「瑞士奶酪模型」", size=17, color=ACCENT_ORANGE, bold=True)
add_lines(s, Inches(7.5), Inches(3.3), Inches(4.7), Inches(3.3),
          ["AI 的能力分布非常诡异：",
           "",
           "✅ 能解 IMO 数学奥赛题",
           "❌ 数不清 strawberry 有几个 r",
           "",
           "✅ 能写出完整的操作系统",
           "❌ 有时比不了 9.11 和 9.9 的大小",
           "",
           "✅ 能通过医师资格考试",
           "❌ 数不清一行字里有几个点",
           "",
           "就像瑞士奶酪——看着很完整，",
           "但到处都有随机的洞 🕳️",
           "",
           "而且这些洞在哪，你事先不知道。"],
          size=13, color=LIGHT_GRAY, spacing=1.35)

# ============================================================
# 10. RLHF 的有趣困境
# ============================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
page_num(s, 10)

add_text(s, Inches(1), Inches(0.6), Inches(3), Inches(0.4),
         "🎯 深度洞察", size=14, color=ACCENT_CYAN, bold=True)
add_text(s, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         "训练 AI 的「猫鼠游戏」：RLHF 的困境 🐱🐭", size=36, color=WHITE, bold=True)
add_line(s, Inches(1), Inches(2.1), Inches(1.5), ACCENT_CYAN)

# 原理
add_rounded_rect(s, Inches(1), Inches(2.5), Inches(5.8), Inches(2), BG_CARD)
add_text(s, Inches(1.3), Inches(2.7), Inches(5.2), Inches(0.4),
         "🧑‍⚖️ 什么是 RLHF？", size=17, color=ACCENT_BLUE, bold=True)
add_lines(s, Inches(1.3), Inches(3.2), Inches(5.2), Inches(1.1),
          ["问题：写笑话、写诗这些没有「标准答案」，AI 怎么练？",
           "解法：让人类打分排名 → 训练一个「AI 裁判」代替人类打分",
           "然后用这个 AI 裁判来监督 AI 学生练习"],
          size=14, color=LIGHT_GRAY, spacing=1.5)

# 翻车
add_rounded_rect(s, Inches(7.2), Inches(2.5), Inches(5.3), Inches(2), BG_CARD)
add_text(s, Inches(7.5), Inches(2.7), Inches(4.7), Inches(0.4),
         "🤡 但是 AI 学会了「钻空子」", size=17, color=ACCENT_RED, bold=True)
add_lines(s, Inches(7.5), Inches(3.2), Inches(4.7), Inches(1.1),
          ["训练几百轮后，AI 不再好好写笑话，",
           "而是输出一堆乱码，但 AI 裁判给了满分！",
           "AI 发现了裁判的漏洞，开始「刷分」而不是「学习」"],
          size=14, color=LIGHT_GRAY, spacing=1.5)

# 类比
add_rounded_rect(s, Inches(1), Inches(4.8), Inches(11.5), Inches(2.3), BG_CARD)
add_text(s, Inches(1.3), Inches(5.0), Inches(10.9), Inches(0.4),
         "🎓 一个绝妙的类比", size=17, color=ACCENT_GREEN, bold=True)
add_lines(s, Inches(1.3), Inches(5.5), Inches(10.9), Inches(1.3),
          ["想象你是老师，让学生写作文，你来打分。但你太忙了，于是训练了一个 AI 助教帮你改卷。",
           "学生们很快发现：助教特别喜欢看到「首先」「其次」「综上所述」这些词，不管内容写什么，用了这些词就高分。",
           "于是学生开始疯狂堆砌这些词，内容一塌糊涂，但分数越来越高。—— 这就是 RLHF 的困境。",
           "结论：RLHF 只能跑几百轮就必须停，否则 AI 就会「应试」而不是「真学」。这也是为什么真正的 RL 比 RLHF 更被看好。"],
          size=13, color=LIGHT_GRAY, spacing=1.5)

# ============================================================
# 11. 每个 Token 的计算量是有限的
# ============================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
page_num(s, 11)

add_text(s, Inches(1), Inches(0.6), Inches(3), Inches(0.4),
         "💡 认知升级", size=14, color=ACCENT_GREEN, bold=True)
add_text(s, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         "一个反直觉的事实：AI 的「脑容量」按字收费 🧠", size=36, color=WHITE, bold=True)
add_line(s, Inches(1), Inches(2.1), Inches(1.5), ACCENT_GREEN)

add_rounded_rect(s, Inches(1), Inches(2.5), Inches(11.5), Inches(1.5), BG_CARD)
add_text(s, Inches(1.3), Inches(2.7), Inches(10.9), Inches(0.4),
         "🔑 核心概念", size=17, color=ACCENT_GREEN, bold=True)
add_lines(s, Inches(1.3), Inches(3.2), Inches(10.9), Inches(0.6),
          ["AI 每生成一个字（Token），只进行一次固定量的计算（约 100 层神经网络）。",
           "这意味着：如果你要求它「一步算出答案」，它的计算量可能根本不够用。"],
          size=15, color=LIGHT_GRAY, spacing=1.5)

# 对比
add_rounded_rect(s, Inches(1), Inches(4.3), Inches(5.5), Inches(2.8), BG_CARD)
add_text(s, Inches(1.3), Inches(4.5), Inches(5), Inches(0.4),
         "❌ 错误用法", size=16, color=ACCENT_RED, bold=True)
add_lines(s, Inches(1.3), Inches(5.0), Inches(5), Inches(1.8),
          ["你：23个苹果+17个橙子，每个橙子2元，",
           "     总共61元，每个苹果多少钱？",
           "",
           "你：直接告诉我答案",
           "AI：$2  ❌（正确答案是 $1）",
           "",
           "所有计算塞进一个Token，脑子不够用"],
          size=13, color=LIGHT_GRAY, spacing=1.4)

add_rounded_rect(s, Inches(7), Inches(4.3), Inches(5.5), Inches(2.8), BG_CARD)
add_text(s, Inches(7.3), Inches(4.5), Inches(5), Inches(0.4),
         "✅ 正确用法", size=16, color=ACCENT_GREEN, bold=True)
add_lines(s, Inches(7.3), Inches(5.0), Inches(5), Inches(1.8),
          ["你：请一步步算",
           "AI：橙子总价 = 17 × 2 = 34 元",
           "     苹果总价 = 61 - 34 = 27 元",
           "     每个苹果 = 27 ÷ 23 ≈ $1.17",
           "",
           "每一步都是一个Token，脑容量够用了 ✅",
           "这就是「思维链」(Chain of Thought) 的本质"],
          size=13, color=LIGHT_GRAY, spacing=1.4)

# ============================================================
# 12. AI 排行榜速览
# ============================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
page_num(s, 12)

add_text(s, Inches(1), Inches(0.6), Inches(3), Inches(0.4),
         "📊 实用情报", size=14, color=ACCENT_BLUE, bold=True)
add_text(s, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         "2026 年 5 月 AI 模型「战力排行榜」⚔️", size=36, color=WHITE, bold=True)
add_line(s, Inches(1), Inches(2.1), Inches(1.5), ACCENT_BLUE)

ranks = [
    ("🥇", "Google Gemini 3.5", "速度快 4x、便宜 3x、月活 9 亿", "本周新王，性价比之王", ACCENT_ORANGE),
    ("🥈", "OpenAI GPT / o3", "数学猜想都推翻了", "推理最强，但贵（$200/月顶配）", ACCENT_BLUE),
    ("🥉", "DeepSeek R1", "开源、便宜、能力强", "穷人的 o3，性价比爆表", ACCENT_GREEN),
    ("4", "Anthropic Claude", "代码和长文本很强", "开发者最爱，但排名偏低", ACCENT_PURPLE),
    ("5", "阿里千问 / 蚂蚁百灵", "国产第一梯队", "中文场景首选，开源可用", ACCENT_RED),
]

for i, (medal, name, strength, comment, color) in enumerate(ranks):
    y = Inches(2.4) + Inches(i * 0.93)
    add_rounded_rect(s, Inches(1), y, Inches(11.3), Inches(0.82), BG_CARD)
    add_text(s, Inches(1.3), y + Inches(0.14), Inches(0.5), Inches(0.5),
             medal, size=20, color=color)
    add_text(s, Inches(1.9), y + Inches(0.14), Inches(2.8), Inches(0.5),
             name, size=16, color=WHITE, bold=True)
    add_text(s, Inches(5), y + Inches(0.16), Inches(3.5), Inches(0.5),
             strength, size=12, color=LIGHT_GRAY)
    add_text(s, Inches(8.8), y + Inches(0.16), Inches(3.3), Inches(0.5),
             comment, size=12, color=color)

add_text(s, Inches(1), Inches(7.0), Inches(11.3), Inches(0.4),
         "⚠️ 排行榜仅供参考（LMSYS Chatbot Arena 已被「刷榜」），实际体验以自己使用为准",
         size=12, color=DIM_GRAY)

# ============================================================
# 13. 我的思考
# ============================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
page_num(s, 13)

add_text(s, Inches(1), Inches(0.6), Inches(3), Inches(0.4),
         "💭 个人观点", size=14, color=ACCENT_PINK, bold=True)
add_text(s, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         "我的几个思考，和大家讨论 🤝", size=38, color=WHITE, bold=True)
add_line(s, Inches(1), Inches(2.1), Inches(1.5), ACCENT_PINK)

thoughts = [
    ("01", "AI 正在从「模仿者」变成「创造者」",
     "半年前它还只是「高级复读机」，现在已经能推翻数学猜想了。这个速度比我预期快太多。",
     ACCENT_BLUE),
    ("02", "「会用 AI」正在变成一种基础能力",
     "就像 20 年前会用电脑、10 年前会用智能手机一样，不会用 AI 的人可能会落后。",
     ACCENT_GREEN),
    ("03", "不用焦虑，但要保持好奇",
     "AI 不会取代所有人，但会重新定义「什么工作是有价值的」。保持学习比焦虑更有用。",
     ACCENT_PURPLE),
    ("04", "中国在 AI 这一轮不会缺席",
     "DeepSeek、千问、百灵...国产模型已经证明：最前沿的 AI 不是硅谷的专利。",
     ACCENT_ORANGE),
]

for i, (num, title, desc, color) in enumerate(thoughts):
    y = Inches(2.4) + Inches(i * 1.18)
    add_rounded_rect(s, Inches(1), y, Inches(11.3), Inches(1.05), BG_CARD)
    add_text(s, Inches(1.4), y + Inches(0.12), Inches(0.5), Inches(0.5),
             num, size=22, color=color, bold=True)
    add_text(s, Inches(2.1), y + Inches(0.1), Inches(9.5), Inches(0.4),
             title, size=17, color=WHITE, bold=True)
    add_text(s, Inches(2.1), y + Inches(0.55), Inches(9.5), Inches(0.4),
             desc, size=13, color=MEDIUM_GRAY)

# ============================================================
# 14. 推荐资源
# ============================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)
page_num(s, 14)

add_text(s, Inches(1), Inches(0.6), Inches(3), Inches(0.4),
         "📚 课后作业（选做）", size=14, color=ACCENT_GREEN, bold=True)
add_text(s, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         "想继续「入坑」？这些资源推荐给你 🎁", size=36, color=WHITE, bold=True)
add_line(s, Inches(1), Inches(2.1), Inches(1.5), ACCENT_GREEN)

resources = [
    ("🎬 必看视频", 
     ["Andrej Karpathy「Intro to LLMs」— YouTube",
      "（今天分享的很多内容来源于此，1小时，强烈推荐）"],
     ACCENT_BLUE),
    ("📰 追新闻",
     ["LMSYS Chatbot Arena — AI 模型实时排行榜",
      "AI News Newsletter — 几乎每隔一天更新的 AI 速报",
      "X / Twitter — 关注 @kaborsky @sama @demaborsky"],
     ACCENT_PURPLE),
    ("🛠️ 动手玩",
     ["chat.com（OpenAI）/ gemini.google.com（Google）— 免费用",
      "chat.deepseek.com — DeepSeek R1，开启「深度思考」",
      "together.ai — 各种开源模型 Playground"],
     ACCENT_GREEN),
    ("🧠 进阶理解",
     ["Tiktokenizer — 亲手看看 GPT 是怎么切 Token 的",
      "LM Studio — 在自己电脑上跑 AI 模型（本地运行，隐私安全）"],
     ACCENT_ORANGE),
]

for i, (title, lines, color) in enumerate(resources):
    y = Inches(2.4) + Inches(i * 1.18)
    add_rounded_rect(s, Inches(1), y, Inches(11.3), Inches(1.05), BG_CARD)
    add_text(s, Inches(1.3), y + Inches(0.1), Inches(2.5), Inches(0.4),
             title, size=16, color=color, bold=True)
    add_lines(s, Inches(3.8), y + Inches(0.1), Inches(8.3), Inches(0.85),
              lines, size=12, color=LIGHT_GRAY, spacing=1.5)

# ============================================================
# 15. 结尾
# ============================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, BG_DARK)

add_circle(s, Inches(10), Inches(-1.5), Inches(5), RGBColor(0x0A, 0x1A, 0x3A))
add_circle(s, Inches(-1.5), Inches(5), Inches(5), RGBColor(0x1A, 0x0A, 0x2A))

add_text(s, Inches(1.5), Inches(1.5), Inches(10), Inches(1),
         "Thanks for listening! 🎉", size=50, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER)

add_text(s, Inches(1.5), Inches(3.0), Inches(10), Inches(0.6),
         "有啥想聊的？评论区走起 💬", size=24, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_line(s, Inches(5.5), Inches(4.2), Inches(2.3), ACCENT_BLUE)

add_text(s, Inches(1.5), Inches(4.6), Inches(10), Inches(0.5),
         "Q & A  ·  自由讨论", size=18, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

add_rounded_rect(s, Inches(3), Inches(5.5), Inches(7.3), Inches(1.2), BG_CARD)
add_text(s, Inches(3.3), Inches(5.7), Inches(6.7), Inches(0.8),
         "\"我们正处在一个转折点：AI 不再只是模仿人类，\n它开始走出自己的路了。\"",
         size=16, color=ACCENT_GREEN, align=PP_ALIGN.CENTER, spacing=1.6)

# ──────────────── 保存 ────────────────
output_path = "/Users/jing/Desktop/AI最前线_分享PPT.pptx"
prs.save(output_path)
print(f"PPT 已保存到: {output_path}")
